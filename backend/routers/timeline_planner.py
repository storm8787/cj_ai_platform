"""
사업 타임라인 생성기 v5
- 4단계: 계획 → 계약 → 시행 → 완료
- 법령 질의 구체화 (사업유형+예산+계약방식별 핀포인트 질의)
- 시행 단계 2회 호출 (1차 큰 공정 → 2차 세부 분해)
- 세부 업무 + 소요기간 + XLSX 시트2
"""

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel, Field
from typing import Optional, List
from datetime import datetime
import json, io, base64, logging, httpx

from services.openai_service import OpenAIService

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api/timeline", tags=["timeline"])
INTERNAL_BASE_URL = "http://localhost:8000"


# ══════════════════════════════════════════════
# Models
# ══════════════════════════════════════════════

class TimelineTask(BaseModel):
    name: str = Field(...)
    start_month: int = Field(..., ge=1, le=12)
    end_month: int = Field(..., ge=1, le=12)
    start_year: int = Field(...)
    end_year: int = Field(...)
    category: Optional[str] = Field(None)
    is_milestone: bool = Field(False)

class TimelineData(BaseModel):
    title: str = Field(...)
    tasks: list[TimelineTask] = Field(...)
    base_year: int = Field(...)

class AutoSuggestRequest(BaseModel):
    project_name: str = Field(...)
    project_description: Optional[str] = None
    budget: Optional[str] = None
    deadline: Optional[str] = None
    project_type: Optional[str] = None
    contract_type: Optional[str] = None

class ExportDetailTask(BaseModel):
    order: Optional[int] = None
    task: str = ""
    description: Optional[str] = None
    legal_basis: Optional[str] = None
    required: bool = False
    note: Optional[str] = None
    duration: Optional[str] = None

class ExportDetailGroup(BaseModel):
    task_name: str = ""
    task_category: str = ""
    items: list[ExportDetailTask] = []

class ExportRequest(BaseModel):
    timeline: TimelineData
    format: str = Field(..., pattern="^(png|xlsx|pptx)$")
    detail_tasks: Optional[list[ExportDetailGroup]] = None

class DetailTasksRequest(BaseModel):
    task_name: str = Field(...)
    task_category: str = Field(...)
    project_name: str = Field(...)
    project_type: Optional[str] = None
    contract_type: Optional[str] = None
    budget: Optional[str] = None
    project_description: Optional[str] = None


# ══════════════════════════════════════════════
# 상수
# ══════════════════════════════════════════════

CATEGORIES = [
    {"value": "계획", "label": "계획", "color": "#7F77DD"},
    {"value": "계약", "label": "계약", "color": "#3B8BD4"},
    {"value": "시행", "label": "시행", "color": "#1D9E75"},
    {"value": "완료", "label": "완료", "color": "#D85A30"},
]

PROJECT_TYPES = [
    {"value": "construction", "label": "건설/토목 공사", "icon": "🏗️"},
    {"value": "it_system", "label": "정보화/시스템 구축", "icon": "💻"},
    {"value": "facility", "label": "시설 설치/개선", "icon": "🏢"},
    {"value": "service", "label": "용역/연구 사업", "icon": "📋"},
    {"value": "event", "label": "행사/축제 기획", "icon": "🎉"},
    {"value": "policy", "label": "정책/제도 개선", "icon": "📝"},
    {"value": "welfare", "label": "복지/지원 사업", "icon": "🤝"},
    {"value": "education", "label": "교육/홍보 사업", "icon": "📚"},
    {"value": "environment", "label": "환경/녹지 사업", "icon": "🌳"},
    {"value": "other", "label": "기타", "icon": "📌"},
]

CONTRACT_TYPES = [
    {"value": "direct_small", "label": "수의계약 (2천만원 이하)", "icon": "📎"},
    {"value": "direct_medium", "label": "소액수의계약 (5천만원 이하)", "icon": "📎"},
    {"value": "limited_qualify", "label": "제한경쟁입찰 (적격심사)", "icon": "📋"},
    {"value": "limited_negotiation", "label": "제한경쟁입찰 (협상에 의한 계약)", "icon": "📋"},
    {"value": "open_competitive", "label": "일반경쟁입찰", "icon": "📢"},
    {"value": "emergency", "label": "긴급계약", "icon": "⚡"},
]

# 사업 유형별 법령 질의 매핑 (구체적 핀포인트 질의)
LAW_QUERIES_BY_TYPE = {
    "it_system": {
        "계획": [
            "전자정부법 정보화사업 보안성 검토 대상 기준",
            "소프트웨어산업진흥법 소프트웨어사업 과업심의 대상 금액",
            "지방자치단체 정보화 사전협의 대상 및 절차",
            "전자정부법 시행령 정보시스템 감리 의무 대상 금액",
            "개인정보보호법 개인정보 영향평가 대상 기준",
        ],
        "계약": [
            "소프트웨어산업진흥법 소프트웨어사업 하도급 제한 기준",
            "소프트웨어 기술성 평가 기준 및 배점",
            "지방계약법 협상에 의한 계약 절차",
        ],
        "완료": [
            "전자정부법 정보시스템 감리 완료 기준",
            "소프트웨어사업 하자보증 기간 및 범위",
            "지방계약법 대가 지급 기한",
        ],
    },
    "construction": {
        "계획": [
            "건설기술진흥법 설계의 경제성 검토 대상 금액",
            "환경영향평가법 소규모 환경영향평가 대상",
            "건설기술진흥법 안전관리계획 수립 대상",
            "건축법 건축심의 대상 기준",
            "국토계획법 개발행위허가 대상",
        ],
        "계약": [
            "지방계약법 적격심사 세부기준",
            "건설산업기본법 건설업 등록 기준",
            "지방계약법 입찰참가자격 사전심사 기준",
        ],
        "완료": [
            "건설기술진흥법 준공검사 절차",
            "건설산업기본법 하자담보책임 기간",
            "지방계약법 계약보증금 및 하자보증금 기준",
            "시설물안전관리특별법 시설물 등록 대상",
        ],
    },
    "service": {
        "계획": [
            "지방자치단체 용역사업 일상감사 대상 금액",
            "지방계약법 수의계약 가능 금액 기준",
        ],
        "계약": [
            "지방계약법 용역계약 일반조건",
            "지방계약법 제안서 평가 기준",
        ],
        "완료": [
            "지방계약법 용역 검수 절차",
            "지방계약법 용역 하자보증 기간",
            "지방계약법 대가 지급 기한",
        ],
    },
    "event": {
        "계획": [
            "지역축제 및 행사 안전관리 계획 수립 기준",
            "도로법 도로점용허가 대상",
            "옥외광고물법 현수막 게시 허가 기준",
            "소음진동관리법 행사장 소음 기준",
        ],
        "계약": [
            "지방계약법 수의계약 가능 금액 기준",
        ],
        "완료": [
            "지방재정법 보조금 정산 기한",
            "지방자치단체 행사 결과보고서 작성 기준",
        ],
    },
    "facility": {
        "계획": [
            "건축법 건축허가 대상 기준",
            "소방시설법 소방시설 설치 기준",
            "장애인편의법 편의시설 설치 대상",
            "에너지이용합리화법 에너지절약계획서 제출 대상",
        ],
        "계약": [
            "지방계약법 적격심사 기준",
            "건설산업기본법 시공 자격 기준",
        ],
        "완료": [
            "건축법 사용승인 절차",
            "건설산업기본법 하자보증 기간",
            "시설물안전관리특별법 시설물 등록",
        ],
    },
}

# 공통 질의 (모든 사업 유형)
COMMON_LAW_QUERIES = {
    "계획": [
        "지방재정법 일상감사 대상 기준 금액",
        "지방재정법 투자심사 대상 기준 금액",
    ],
    "계약": [
        "지방계약법 계약보증금 납부 기준",
    ],
    "완료": [
        "지방계약법 대가 지급 기한",
        "지방계약법 하자보증금 기준",
    ],
}

# 예산 규모별 추가 질의
BUDGET_LAW_QUERIES = [
    {"threshold": "1억", "query": "지방재정법 1억원 이상 사업 투자심사 기준"},
    {"threshold": "5억", "query": "지방재정법 5억원 이상 사업 타당성조사 기준"},
    {"threshold": "10억", "query": "건설기술진흥법 10억원 이상 건설공사 안전관리계획"},
]


def _clean_json(text):
    """GPT 응답에서 JSON 추출"""
    c = text.strip()
    if c.startswith("```"):
        c = c.split("\n", 1)[1] if "\n" in c else c
    if c.endswith("```"):
        c = c[:-3]
    return c.strip()


def _tl(v):
    """사업 유형 value → label"""
    return next((t["label"] for t in PROJECT_TYPES if t["value"] == v), v) if v else ""


def _cl(v):
    """계약 방식 value → label"""
    return next((c["label"] for c in CONTRACT_TYPES if c["value"] == v), v) if v else ""


# ══════════════════════════════════════════════
# 법령 챗봇 연동
# ══════════════════════════════════════════════

async def _ask_law(question: str) -> str:
    """법령 챗봇 API 내부 호출"""
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            resp = await client.post(
                f"{INTERNAL_BASE_URL}/api/law-chatbot/ask",
                json={"question": question, "search_scope": "all"}
            )
            if resp.status_code == 200:
                return resp.json().get("answer", "")
    except Exception as e:
        logger.warning(f"법령 챗봇 연동 실패: {e}")
    return ""


def _build_law_queries(project_type: str, category: str, budget: str = None) -> list:
    """사업유형 + 카테고리 + 예산에 따른 구체적 법령 질의 목록 생성"""
    queries = []

    # 사업 유형별 질의
    type_queries = LAW_QUERIES_BY_TYPE.get(project_type, {}).get(category, [])
    queries.extend(type_queries)

    # 공통 질의
    common = COMMON_LAW_QUERIES.get(category, [])
    queries.extend(common)

    # 예산 규모별 추가 질의
    if budget and category == "계획":
        for bq in BUDGET_LAW_QUERIES:
            queries.append(bq["query"])

    # 사업 유형이 매핑에 없는 경우 일반적 질의
    if not type_queries:
        general_queries = {
            "계획": [f"지방자치단체 사업 사전 행정절차 심의 검토 협의 기준"],
            "계약": [f"지방계약법 계약 체결 절차"],
            "완료": [f"지방계약법 사업 완료 검수 정산 절차"],
        }
        queries.extend(general_queries.get(category, []))

    return queries


async def _get_law_context(project_type: str, category: str, budget: str = None) -> str:
    """법령 질의를 실행하고 결과를 컨텍스트 문자열로 반환"""
    queries = _build_law_queries(project_type, category, budget)
    results = []

    for q in queries:
        answer = await _ask_law(q)
        if answer and len(answer) > 20:  # 너무 짧은 답변은 무시
            results.append(f"[질의: {q}]\n{answer}")

    if results:
        return "\n\n---\n\n".join(results)
    return ""


# ══════════════════════════════════════════════
# 일정 추천
# ══════════════════════════════════════════════

SUGGEST_PROMPT = """당신은 한국 지방자치단체의 사업 일정 전문가입니다.

사업 일정을 4단계로 구분하여 추천하세요:
- "계획": 기본계획 수립, 사전 심의/검토, 일상감사, 예산확보 등
- "계약": 설계서/과업지시서 작성, 입찰공고, 제안평가, 계약체결 등
- "시행": 실제 사업 수행 (공사, 개발, 용역수행 등)
- "완료": 준공검사/검수, 대가지급, 정산, 하자보증 등

규칙:
1. category는 반드시 "계획", "계약", "시행", "완료" 중 하나
2. is_milestone은 항상 false
3. 완료 목표 시기 지정 시 반드시 준수. 초과 불가. 부족하면 압축.
4. 사업 설명의 기술/방법론을 단계명과 summary에 반영
5. 계약 방식 지정 시 해당 절차 반영
6. 상하반기 인사이동(1,7월) 주요 일정 시작 지양
7. 연말(11~12월) 결산 고려
8. 아래 소요기간 기준을 반드시 참고하여 현실적으로 일정을 산출하세요. 비현실적으로 늘리거나 줄이지 마세요.

[단계별 현실적 소요기간 기준 — 지방계약법 기반]

■ 계획 단계
  - 기본계획 수립: 2~4주
  - 사전 심의/검토 (보안성검토, 과업심의, 설계심의 등): 각 1~2주
  - 일상감사: 3~7일
  - 투자심사 (5억 이상): 2~4주
  - 환경영향평가 (해당 시): 2~3개월

■ 계약 단계
  - 설계서/과업지시서/제안요청서 작성: 2~4주
  - 원가계산/예정가격 산정: 1~2주
  - 입찰공고 기간 (지방계약법 시행령 제35조):
    · 일반입찰: 7일
    · 협상에 의한 계약: 추정가격 1억 미만 10일, 1억~10억 20일, 10억 이상 40일
    · 긴급/재공고: 5일
  - 현장설명 (공사, 해당 시): 공고 후 7일 이후
  - 제안서 평가/기술평가위원회: 1~2주
  - 적격심사: 서류제출 후 7일 이내 (불가피 시 +3일)
  - 낙찰자 결정 통보: 심사 완료 후 즉시~3일
  - 계약 체결 (계약서 작성, 보증보험 징구): 낙찰 후 7~10일
  ※ 수의계약: 견적 징구~계약 체결까지 1~2주
  ※ 소액수의계약: 견적비교~계약 체결까지 1~2주

■ 시행 단계
  - 사업 규모와 유형에 따라 산정 (총 3~12개월)
  - 착수신고/착수보고: 계약 후 7~14일
  - 중간보고 (용역): 시행기간의 약 50% 시점

■ 완료 단계
  - 준공검사/검수: 1~2주
  - 대가 지급: 검수 후 14일 이내 (지방계약법 제19조)
  - 정산/결산 서류 작성: 1~2주
  - 하자보증보험 징구: 계약 종료 시점

반드시 아래 JSON만 반환:
{
  "tasks": [{"name": "단계명", "start_month": 3, "end_month": 4, "start_year": 2026, "end_year": 2026, "category": "계획", "is_milestone": false}],
  "summary": "일정 산출 근거 2~3문장"
}"""


@router.post("/suggest")
async def suggest_timeline(request: AutoSuggestRequest):
    """GPT 기반 자동 일정 추천"""
    try:
        svc = OpenAIService()
        parts = [f"사업명: {request.project_name}"]
        if request.project_description:
            parts.append(f"사업 설명: {request.project_description}")
        if request.budget:
            parts.append(f"예산 규모: {request.budget}")
        if request.deadline:
            parts.append(f"완료 목표: {request.deadline}")
        if request.project_type:
            parts.append(f"사업 유형: {_tl(request.project_type)}")
        if request.contract_type:
            parts.append(f"계약 방식: {_cl(request.contract_type)}")

        now = datetime.now()
        parts.append(f"현재 시점: {now.year}년 {now.month}월")
        parts.append("위 사업의 현실적인 추진 일정을 추천해 주세요.")

        result_text = await svc.generate_text(
            prompt=f"{SUGGEST_PROMPT}\n\n" + "\n".join(parts),
            max_tokens=2000,
            temperature=0.7
        )
        result = json.loads(_clean_json(result_text))

        valid_cats = {"계획", "계약", "시행", "완료"}
        validated = []
        for t in result.get("tasks", []):
            cat = t.get("category", "시행")
            if cat not in valid_cats:
                cat = "시행"
            validated.append({
                "name": t.get("name", "미정"),
                "start_month": max(1, min(12, t.get("start_month", 1))),
                "end_month": max(1, min(12, t.get("end_month", 1))),
                "start_year": t.get("start_year", now.year),
                "end_year": t.get("end_year", now.year),
                "category": cat,
                "is_milestone": False,
            })

        return {
            "success": True,
            "tasks": validated,
            "summary": result.get("summary", ""),
            "project_name": request.project_name,
        }

    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail="AI 응답을 파싱할 수 없습니다.")
    except Exception as e:
        logger.error(f"일정 추천 실패: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ══════════════════════════════════════════════
# 세부 업무 생성 (단계별 분기)
# ══════════════════════════════════════════════

# 계획/계약 단계: 법령 연동 + 구체적 세부업무
DETAIL_WITH_LAW = """당신은 한국 지방자치단체의 사업 관리 전문가입니다.
주어진 단계의 세부 업무 목록을 생성하세요.

규칙:
1. 아래 법령 검색 결과를 기반으로 해당 사업에 필요한 사전 절차를 구체적으로 나열하세요.
2. 각 업무의 예상 소요기간(duration)을 현실적으로 산출하세요.
3. 법령에서 확인된 의무사항은 required: true, legal_basis에 근거 법령을 명시하세요.
4. 예산 규모에 따라 해당/비해당을 판단하여 해당되는 것만 포함하세요.
5. 사업 유형과 계약 방식에 맞는 절차만 선별하세요.
6. 단순히 "검토", "협의" 등 뭉뚱그리지 말고, 구체적으로 어떤 검토인지, 어디에 협의하는지 명시하세요.
7. 순서는 실제 진행 순서대로 나열하세요.
8. description은 반드시 3문장 이상으로 상세하게 작성하세요:
   - 누가 (담당부서/담당자)
   - 무엇을 (구체적인 서류명, 시스템명, 절차명)
   - 어떻게 (제출처, 심의기관, 처리방법, 필요서류)
   예시 (좋은 description):
   "정보화담당관실에 정보화사업 사전협의서를 작성하여 제출한다. 협의서에는 사업 개요, 시스템 구성도, 기존 시스템 연계방안, 개인정보처리 여부를 포함해야 한다. 협의 결과에 따라 사업계획을 보완한 후 다음 단계로 진행한다."
   예시 (나쁜 description):
   "사전협의를 진행한다."

반드시 아래 JSON만 반환:
{
  "detail_tasks": [
    {"order": 1, "task": "구체적 업무명", "description": "누가, 무엇을, 어떻게 하는지 3문장 이상의 상세 설명", "duration": "약 1주", "legal_basis": "근거법령 제X조", "required": true, "note": "참고사항 또는 null"}
  ]
}"""


# 시행 단계 1차: 큰 공정 분해
DETAIL_EXECUTE_PHASE1 = """당신은 한국 지방자치단체의 사업 관리 전문가입니다.
사업의 시행 단계를 세부 작업 공정으로 분해하세요.

규칙:
1. 사업 설명에 언급된 내용을 기반으로 실제 수행할 세부 작업을 구체적으로 나열하세요.
2. 각 작업의 예상 소요기간을 현실적으로 산출하세요.
3. 아래 사업 유형별 세부 공정 예시를 참고하되, 사업 설명에 맞게 구체화하세요.

[정보화/시스템 구축 예시]
- 요구사항 분석 및 정의 (업무 현황 조사, 사용자 인터뷰, 요구사항 명세서 작성)
- 시스템 설계 (아키텍처 설계, DB 설계, UI/UX 설계, 인터페이스 설계)
- 인프라 구축 (서버 도입 및 설치, 네트워크 구성, 보안장비 설치)
- 소프트웨어 개발 (프론트엔드, 백엔드, API, 배치 프로그램)
- 단위 테스트 (모듈별 기능 테스트, 버그 수정)
- 통합 테스트 (시스템 간 연계 테스트, 성능 테스트, 보안 테스트)
- 데이터 이관 (기존 데이터 정제, 매핑, 이관, 검증)
- 시범 운영 (실사용자 테스트, 오류 수정, 매뉴얼 작성)
- 사용자 교육 (관리자 교육, 실무자 교육, 교육자료 제작)

[건설/토목 공사 예시]
- 현장 준비 (가설 울타리 설치, 가설 사무소, 교통안전시설)
- 토공 (터파기, 성토, 다짐, 잔토 처리)
- 기초 공사 (기초 터파기, 철근 배근, 콘크리트 타설)
- 골조 공사 (철근콘크리트, 철골, 조적)
- 포장 공사 (노반 정리, 기층, 표층, 아스콘 포설)
- 부대 공사 (배수관, 맨홀, 경계석, 가드레일)
- 조경 공사 (식재, 잔디, 시설물)
- 마감 및 정리 (가설물 철거, 현장 정리, 도로 청소)

[용역/연구 사업 예시]
- 착수 보고 (연구 계획 발표, 자문위원 구성)
- 현황 조사 (문헌 조사, 현장 조사, 설문 조사, 벤치마킹)
- 분석 (데이터 분석, 문제점 도출, 시사점 정리)
- 중간 보고 (중간 성과 발표, 자문위원 검토, 방향 조정)
- 대안 수립 (대안 도출, 비교 분석, 최적안 선정)
- 성과물 작성 (보고서 작성, 도면 작성, 매뉴얼 작성)
- 최종 보고 (최종 성과 발표, 자문위원 심의, 수정 보완)

[행사/축제 예시]
- 세부 기획 (프로그램 구성, 동선 설계, 예산 배분)
- 출연진/업체 섭외 (공연팀, 부스운영, 장비업체)
- 홍보 (포스터, SNS, 보도자료, 현수막)
- 시설 설치 (무대, 음향, 조명, 부스, 안내판)
- 리허설 (동선 점검, 음향 테스트, 안전 점검)
- 본 행사 운영 (진행, 안전관리, 운영인력 배치)
- 철거 및 정리 (시설 철거, 현장 정리, 쓰레기 수거)

4. "개발" 이라고 한 줄로 끝내지 말고, 실제로 무엇을 개발하는지 (프론트엔드, 백엔드, DB, API 등) 구체적으로 나누세요.
5. "공사 진행" 이라고 한 줄로 끝내지 말고, 어떤 공종을 어떤 순서로 하는지 나누세요.
6. description은 반드시 3문장 이상으로 상세하게 작성하세요. 사업 설명에서 언급된 구체적인 기술, 장비, 시스템을 반영하세요.
   예시 (좋은 description):
   "LLM 기반 챗봇 백엔드를 FastAPI로 개발하며, OpenAI GPT-4o API 연동과 RAG 파이프라인(FAISS 벡터스토어 + BM25 하이브리드 검색)을 구축한다. 법령 데이터 12,000건의 임베딩 처리 및 청크 분할 로직을 구현하고, 답변 생성 시 환각 방지를 위한 프롬프트 엔지니어링을 적용한다. API 응답 속도 최적화를 위해 캐싱 레이어와 비동기 처리를 도입한다."
   예시 (나쁜 description):
   "백엔드를 개발한다."

반드시 아래 JSON만 반환:
{
  "detail_tasks": [
    {"order": 1, "task": "구체적 작업명", "description": "무엇을 어떻게 하는지 3문장 이상의 상세 설명 (사업 설명에 언급된 기술/장비/시스템 반영)", "duration": "약 2주", "legal_basis": null, "required": true, "note": "참고사항 또는 null"}
  ]
}"""


# 완료 단계: 법령 + 사업유형별 마무리
DETAIL_COMPLETE = """당신은 한국 지방자치단체의 사업 관리 전문가입니다.
사업 완료 단계의 세부 업무를 생성하세요.

규칙:
1. 법령 검색 결과를 기반으로 법정 필수 절차를 먼저 나열하세요.
2. 각 업무의 예상 소요기간을 산출하세요.
3. 사업 유형별 마무리 업무를 포함하세요:
   - 정보화: 데이터 이관 완료 확인, 운영 인수인계서 작성, 관리자/실무자 교육, 유지보수 계약 체결, 소스코드/산출물 납품
   - 건설: 준공도서 작성, 시설물 등록, 관리 부서 이관, 하자보증보험 징구
   - 용역: 최종보고회 개최, 성과심의위원회 심의, 성과물 납품, 저작권 처리
   - 행사: 정산서 작성, 결과보고서 작성, 성과분석(참여인원, 만족도)
4. 법적 근거가 있는 항목은 근거 조문까지 명시하세요.
5. 순서는 실제 진행 순서대로 나열하세요.
6. description은 반드시 3문장 이상으로 상세하게 작성하세요:
   - 구체적으로 어떤 서류를 작성하는지
   - 어디에 제출/보고하는지
   - 기한이나 주의사항은 무엇인지
   예시 (좋은 description):
   "검수조서를 작성하여 계약 상대방에게 통보하고, 납품된 시스템이 과업지시서의 요구사항을 충족하는지 기능별 체크리스트를 기반으로 확인한다. 검수 시 하자가 발견되면 보완 요청서를 발급하고 보완 완료 후 재검수를 실시한다. 검수 완료 후 14일 이내에 대가를 지급해야 한다."
   예시 (나쁜 description):
   "검수를 진행한다."

반드시 아래 JSON만 반환:
{
  "detail_tasks": [
    {"order": 1, "task": "업무명", "description": "누가, 무엇을, 어떻게 하는지 3문장 이상의 상세 설명", "duration": "약 3일", "legal_basis": "근거법령 제X조", "required": true, "note": "참고사항 또는 null"}
  ]
}"""


@router.post("/detail-tasks")
async def generate_detail_tasks(request: DetailTasksRequest):
    """단계별 세부 업무 자동 생성"""
    try:
        svc = OpenAIService()
        tl = _tl(request.project_type)
        cl = _cl(request.contract_type)
        cat = request.task_category
        law_ctx = ""

        # ── 계획 단계: 법령 핀포인트 질의 ──
        if cat == "계획":
            law_ctx = await _get_law_context(
                request.project_type or "", cat, request.budget
            )
            base_prompt = DETAIL_WITH_LAW

        # ── 계약 단계: 법령 핀포인트 질의 ──
        elif cat == "계약":
            law_ctx = await _get_law_context(
                request.project_type or "", cat, request.budget
            )
            base_prompt = DETAIL_WITH_LAW

        # ── 시행 단계: GPT 2회 호출 (세부 분해) ──
        elif cat == "시행":
            base_prompt = DETAIL_EXECUTE_PHASE1

        # ── 완료 단계: 법령 + GPT 혼합 ──
        elif cat == "완료":
            law_ctx = await _get_law_context(
                request.project_type or "", cat, request.budget
            )
            base_prompt = DETAIL_COMPLETE

        else:
            base_prompt = DETAIL_WITH_LAW

        # 사용자 정보 구성
        user_parts = [
            f"사업명: {request.project_name}",
            f"현재 단계: {request.task_name} ({cat})",
        ]
        if tl:
            user_parts.append(f"사업 유형: {tl}")
        if cl:
            user_parts.append(f"계약 방식: {cl}")
        if request.budget:
            user_parts.append(f"예산 규모: {request.budget}")
        if request.project_description:
            user_parts.append(f"사업 설명: {request.project_description}")
        user_parts.append("\n위 단계의 세부 업무와 각 업무별 예상 소요기간을 생성해 주세요.")

        # 프롬프트 조합
        if law_ctx:
            full_prompt = f"{base_prompt}\n\n[법령 검색 결과]\n{law_ctx}\n\n" + "\n".join(user_parts)
        else:
            full_prompt = f"{base_prompt}\n\n" + "\n".join(user_parts)

        # GPT 호출
        result_text = await svc.generate_text(
            prompt=full_prompt,
            max_tokens=3000,
            temperature=0.5
        )
        result = json.loads(_clean_json(result_text))
        detail_tasks = result.get("detail_tasks", [])

        # 시행 단계: 항목이 5개 이하면 2차 세부 분해 시도
        if cat == "시행" and len(detail_tasks) <= 5 and request.project_description:
            task_names = [dt.get("task", "") for dt in detail_tasks]
            refine_prompt = f"""아래는 "{request.project_name}" 사업의 시행 단계 1차 분해 결과입니다.
각 작업을 더 구체적인 세부 작업으로 분해하고, 소요기간을 산출해 주세요.

사업 설명: {request.project_description}
사업 유형: {tl}
1차 분해 결과: {', '.join(task_names)}

각 1차 작업을 2~4개의 세부 작업으로 더 나눠주세요.
전체를 하나의 detail_tasks 배열로 통합하여 순서대로 반환하세요.

반드시 아래 JSON만 반환:
{{"detail_tasks": [{{"order": 1, "task": "세부 작업명", "description": "상세 설명 (2~3문장)", "duration": "약 1주", "legal_basis": null, "required": true, "note": null}}]}}"""

            try:
                refine_text = await svc.generate_text(
                    prompt=refine_prompt,
                    max_tokens=3000,
                    temperature=0.5
                )
                refine_result = json.loads(_clean_json(refine_text))
                refined_tasks = refine_result.get("detail_tasks", [])
                if len(refined_tasks) > len(detail_tasks):
                    detail_tasks = refined_tasks
            except Exception as e:
                logger.warning(f"시행 단계 2차 분해 실패 (1차 결과 사용): {e}")

        return {
            "success": True,
            "task_name": request.task_name,
            "task_category": cat,
            "detail_tasks": detail_tasks,
            "law_referenced": bool(law_ctx),
        }

    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail="AI 응답을 파싱할 수 없습니다.")
    except Exception as e:
        logger.error(f"세부 업무 생성 실패: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ══════════════════════════════════════════════
# 목록 조회
# ══════════════════════════════════════════════

@router.get("/project-types")
async def get_project_types():
    return {"types": PROJECT_TYPES}

@router.get("/contract-types")
async def get_contract_types():
    return {"types": CONTRACT_TYPES}

@router.get("/categories")
async def get_categories():
    return {"categories": CATEGORIES}


# ══════════════════════════════════════════════
# PNG 내보내기
# ══════════════════════════════════════════════

def _generate_png(timeline: TimelineData) -> bytes:
    from PIL import Image, ImageDraw, ImageFont
    import os

    tasks = timeline.tasks
    title = timeline.title
    LEFT_LABEL_W = 280
    MONTH_COL_W = 100
    ROW_H = 56
    HEADER_H = 70
    TITLE_H = 50
    LEGEND_H = 50
    PADDING = 24

    all_months = []
    for t in tasks:
        for y in range(t.start_year, t.end_year + 1):
            sm = t.start_month if y == t.start_year else 1
            em = t.end_month if y == t.end_year else 12
            for m in range(sm, em + 1):
                all_months.append((y, m))
    if not all_months:
        all_months = [(timeline.base_year, m) for m in range(1, 13)]
    all_months = sorted(set(all_months))
    num_months = len(all_months)
    month_index = {ym: i for i, ym in enumerate(all_months)}

    CHART_W = num_months * MONTH_COL_W
    TOTAL_W = LEFT_LABEL_W + CHART_W + PADDING * 2
    TOTAL_H = TITLE_H + HEADER_H + len(tasks) * ROW_H + LEGEND_H + PADDING * 2

    img = Image.new("RGB", (TOTAL_W, TOTAL_H), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    font_paths = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    font = font_bold = font_small = None
    for fp in font_paths:
        if os.path.exists(fp):
            try:
                font = ImageFont.truetype(fp, 16)
                font_bold = ImageFont.truetype(fp, 18)
                font_small = ImageFont.truetype(fp, 13)
                break
            except Exception:
                continue
    if font is None:
        font = ImageFont.load_default()
        font_bold = font
        font_small = font

    CAT_COLORS = {
        "계획": {"fill": (238, 237, 254), "bar": (127, 119, 221), "text": (60, 52, 137)},
        "계약": {"fill": (230, 241, 251), "bar": (59, 139, 212), "text": (12, 68, 124)},
        "시행": {"fill": (225, 245, 238), "bar": (29, 158, 117), "text": (8, 80, 65)},
        "완료": {"fill": (250, 236, 231), "bar": (216, 90, 48), "text": (113, 43, 19)},
    }
    DEFAULT_COLOR = {"fill": (240, 240, 240), "bar": (150, 150, 150), "text": (80, 80, 80)}

    bbox = draw.textbbox((0, 0), title, font=font_bold)
    tw = bbox[2] - bbox[0]
    draw.text(((TOTAL_W - tw) / 2, 16), title, fill=(44, 44, 42), font=font_bold)

    chart_x = PADDING + LEFT_LABEL_W
    header_y = TITLE_H

    for i, (year, month) in enumerate(all_months):
        x = chart_x + i * MONTH_COL_W
        label = f"{year}년 {month}월" if (month == 1 or i == 0) else f"{month}월"
        draw.rectangle([x, header_y, x + MONTH_COL_W, header_y + 30], fill=(241, 239, 232))
        draw.rectangle([x, header_y, x + MONTH_COL_W, header_y + 30], outline=(211, 209, 199))
        bbox = draw.textbbox((0, 0), label, font=font_small)
        mw = bbox[2] - bbox[0]
        draw.text((x + (MONTH_COL_W - mw) / 2, header_y + 8), label, fill=(68, 68, 65), font=font_small)

    row_top = header_y + 30
    for idx, task in enumerate(tasks):
        y = row_top + idx * ROW_H
        colors = CAT_COLORS.get(task.category, DEFAULT_COLOR)
        if idx % 2 == 0:
            draw.rectangle([PADDING, y, TOTAL_W - PADDING, y + ROW_H], fill=(250, 250, 248))
        draw.line([PADDING, y + ROW_H, TOTAL_W - PADDING, y + ROW_H], fill=(230, 228, 222), width=1)
        draw.text((PADDING + 12, y + (ROW_H - 18) / 2), task.name, fill=(44, 44, 42), font=font)

        start_idx = month_index.get((task.start_year, task.start_month), 0)
        end_idx = month_index.get((task.end_year, task.end_month), num_months - 1)
        bar_x1 = chart_x + start_idx * MONTH_COL_W + 6
        bar_x2 = chart_x + (end_idx + 1) * MONTH_COL_W - 6
        bar_y = y + 14
        bar_h = ROW_H - 28

        draw.rounded_rectangle([bar_x1, bar_y, bar_x2, bar_y + bar_h], radius=5,
                               fill=colors["fill"], outline=colors["bar"], width=2)
        months_span = end_idx - start_idx + 1
        if months_span > 1 and (bar_x2 - bar_x1) > 60:
            span_text = f"{months_span}개월"
            bbox = draw.textbbox((0, 0), span_text, font=font_small)
            stw = bbox[2] - bbox[0]
            draw.text(((bar_x1 + bar_x2 - stw) / 2, bar_y + 4), span_text,
                      fill=colors["text"], font=font_small)

    for i in range(num_months + 1):
        x = chart_x + i * MONTH_COL_W
        draw.line([x, row_top, x, row_top + len(tasks) * ROW_H], fill=(238, 236, 230), width=1)

    legend_y = row_top + len(tasks) * ROW_H + 16
    lx = PADDING + 16
    for cat_name, colors in CAT_COLORS.items():
        draw.rounded_rectangle([lx, legend_y + 2, lx + 14, legend_y + 16], radius=3, fill=colors["bar"])
        draw.text((lx + 20, legend_y), cat_name, fill=(68, 68, 65), font=font_small)
        bbox = draw.textbbox((0, 0), cat_name, font=font_small)
        lx += 20 + (bbox[2] - bbox[0]) + 28

    draw.rectangle([0, 0, TOTAL_W - 1, TOTAL_H - 1], outline=(211, 209, 199), width=1)

    buf = io.BytesIO()
    img.save(buf, format="PNG", dpi=(150, 150))
    return buf.getvalue()


# ══════════════════════════════════════════════
# XLSX 내보내기 (시트2: 세부업무+소요기간)
# ══════════════════════════════════════════════

def _generate_xlsx(timeline: TimelineData, detail_groups=None) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()
    ws = wb.active
    ws.title = "사업추진일정"

    all_months = []
    for t in timeline.tasks:
        for y in range(t.start_year, t.end_year + 1):
            sm = t.start_month if y == t.start_year else 1
            em = t.end_month if y == t.end_year else 12
            for m in range(sm, em + 1):
                all_months.append((y, m))
    all_months = sorted(set(all_months))
    month_index = {ym: i for i, ym in enumerate(all_months)}

    header_font = Font(name="맑은 고딕", size=14, bold=True)
    col_header_font = Font(name="맑은 고딕", size=10, bold=True, color="FFFFFF")
    cell_font = Font(name="맑은 고딕", size=10)
    header_fill = PatternFill(start_color="2D3748", end_color="2D3748", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin", color="D1D5DB"),
        right=Side(style="thin", color="D1D5DB"),
        top=Side(style="thin", color="D1D5DB"),
        bottom=Side(style="thin", color="D1D5DB"),
    )

    CATEGORY_FILLS = {
        "계획": PatternFill(start_color="EEEDFE", end_color="EEEDFE", fill_type="solid"),
        "계약": PatternFill(start_color="E6F1FB", end_color="E6F1FB", fill_type="solid"),
        "시행": PatternFill(start_color="E1F5EE", end_color="E1F5EE", fill_type="solid"),
        "완료": PatternFill(start_color="FAECE7", end_color="FAECE7", fill_type="solid"),
    }
    default_fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")

    # 시트1: 간트차트
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=3 + len(all_months))
    title_cell = ws.cell(row=1, column=1, value=timeline.title)
    title_cell.font = header_font
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 36

    headers = ["단계명", "시작", "종료"] + [
        f"{y}년 {m}월" if m == 1 or i == 0 else f"{m}월"
        for i, (y, m) in enumerate(all_months)
    ]
    for col_idx, header in enumerate(headers, 1):
        cell = ws.cell(row=3, column=col_idx, value=header)
        cell.font = col_header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = thin_border

    ws.column_dimensions["A"].width = 30
    ws.column_dimensions["B"].width = 12
    ws.column_dimensions["C"].width = 12
    for i in range(len(all_months)):
        col_letter = chr(68 + i) if i < 22 else None
        if col_letter:
            ws.column_dimensions[col_letter].width = 6

    for row_idx, task in enumerate(timeline.tasks, 4):
        ws.cell(row=row_idx, column=1, value=task.name).font = cell_font
        ws.cell(row=row_idx, column=1).border = thin_border
        ws.cell(row=row_idx, column=2, value=f"{task.start_year}.{task.start_month:02d}").font = cell_font
        ws.cell(row=row_idx, column=2).alignment = Alignment(horizontal="center")
        ws.cell(row=row_idx, column=2).border = thin_border
        ws.cell(row=row_idx, column=3, value=f"{task.end_year}.{task.end_month:02d}").font = cell_font
        ws.cell(row=row_idx, column=3).alignment = Alignment(horizontal="center")
        ws.cell(row=row_idx, column=3).border = thin_border

        start_idx = month_index.get((task.start_year, task.start_month), 0)
        end_idx = month_index.get((task.end_year, task.end_month), len(all_months) - 1)
        fill = CATEGORY_FILLS.get(task.category, default_fill)

        for i in range(len(all_months)):
            cell = ws.cell(row=row_idx, column=4 + i)
            cell.border = thin_border
            if start_idx <= i <= end_idx:
                cell.fill = fill
                cell.value = "■"
                cell.alignment = Alignment(horizontal="center")
                cell.font = Font(size=10, color="666666")

    # 시트2: 세부 업무 (TODO + 소요기간)
    if detail_groups:
        ws2 = wb.create_sheet("세부업무(TODO)")
        ws2.merge_cells(start_row=1, start_column=1, end_row=1, end_column=9)
        title_cell2 = ws2.cell(row=1, column=1, value=f"{timeline.title} — 세부 업무 목록")
        title_cell2.font = header_font
        title_cell2.alignment = Alignment(horizontal="center", vertical="center")
        ws2.row_dimensions[1].height = 36

        todo_headers = ["단계", "구분", "순서", "업무명", "소요기간", "설명", "법적 근거", "필수", "참고사항"]
        for col_idx, header in enumerate(todo_headers, 1):
            cell = ws2.cell(row=3, column=col_idx, value=header)
            cell.font = col_header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border

        ws2.column_dimensions["A"].width = 20
        ws2.column_dimensions["B"].width = 8
        ws2.column_dimensions["C"].width = 6
        ws2.column_dimensions["D"].width = 30
        ws2.column_dimensions["E"].width = 12
        ws2.column_dimensions["F"].width = 50
        ws2.column_dimensions["G"].width = 30
        ws2.column_dimensions["H"].width = 6
        ws2.column_dimensions["I"].width = 30

        row = 4
        for grp in detail_groups:
            cat_fill = CATEGORY_FILLS.get(grp.task_category, default_fill)
            for dt in grp.items:
                ws2.cell(row=row, column=1, value=grp.task_name).font = cell_font
                ws2.cell(row=row, column=1).border = thin_border
                ws2.cell(row=row, column=1).fill = cat_fill

                ws2.cell(row=row, column=2, value=grp.task_category).font = cell_font
                ws2.cell(row=row, column=2).border = thin_border
                ws2.cell(row=row, column=2).alignment = Alignment(horizontal="center")

                ws2.cell(row=row, column=3, value=dt.order or 0).font = cell_font
                ws2.cell(row=row, column=3).border = thin_border
                ws2.cell(row=row, column=3).alignment = Alignment(horizontal="center")

                ws2.cell(row=row, column=4, value=dt.task).font = cell_font
                ws2.cell(row=row, column=4).border = thin_border

                ws2.cell(row=row, column=5, value=dt.duration or "").font = cell_font
                ws2.cell(row=row, column=5).border = thin_border
                ws2.cell(row=row, column=5).alignment = Alignment(horizontal="center")

                ws2.cell(row=row, column=6, value=dt.description or "").font = cell_font
                ws2.cell(row=row, column=6).border = thin_border

                ws2.cell(row=row, column=7, value=dt.legal_basis or "").font = cell_font
                ws2.cell(row=row, column=7).border = thin_border

                ws2.cell(row=row, column=8, value="O" if dt.required else "").font = cell_font
                ws2.cell(row=row, column=8).border = thin_border
                ws2.cell(row=row, column=8).alignment = Alignment(horizontal="center")

                ws2.cell(row=row, column=9, value=dt.note or "").font = cell_font
                ws2.cell(row=row, column=9).border = thin_border

                row += 1

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════
# PPTX 내보내기
# ══════════════════════════════════════════════

def _generate_pptx(timeline: TimelineData) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = timeline.title
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0x2D, 0x37, 0x48)
    title_p.alignment = PP_ALIGN.LEFT

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.4))
    subtitle_p = subtitle_box.text_frame.paragraphs[0]
    subtitle_p.text = "사업 추진 일정표"
    subtitle_p.font.size = Pt(14)
    subtitle_p.font.color.rgb = RGBColor(0x71, 0x71, 0x71)

    # 월 범위 계산
    all_months = []
    for t in timeline.tasks:
        for y in range(t.start_year, t.end_year + 1):
            sm = t.start_month if y == t.start_year else 1
            em = t.end_month if y == t.end_year else 12
            for m in range(sm, em + 1):
                all_months.append((y, m))
    all_months = sorted(set(all_months))
    month_index = {ym: i for i, ym in enumerate(all_months)}
    num_months = len(all_months)

    LEFT = Inches(0.5)
    TOP = Inches(1.6)
    LABEL_W = Inches(2.8)
    CHART_W = Inches(9.5)
    ROW_H = Inches(0.45)
    MONTH_W = CHART_W / num_months if num_months > 0 else Inches(1)

    CAT_COLORS = {
        "계획": RGBColor(0x7F, 0x77, 0xDD),
        "계약": RGBColor(0x3B, 0x8B, 0xD4),
        "시행": RGBColor(0x1D, 0x9E, 0x75),
        "완료": RGBColor(0xD8, 0x5A, 0x30),
    }

    # 월 헤더
    for i, (year, month) in enumerate(all_months):
        x = LEFT + LABEL_W + int(MONTH_W * i)
        label = f"'{str(year)[2:]}.{month}월" if (month == 1 or i == 0) else f"{month}월"
        header_box = slide.shapes.add_textbox(x, TOP - Inches(0.35), int(MONTH_W), Inches(0.3))
        header_p = header_box.text_frame.paragraphs[0]
        header_p.text = label
        header_p.font.size = Pt(9)
        header_p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        header_p.alignment = PP_ALIGN.CENTER

    # 행
    for idx, task in enumerate(timeline.tasks):
        y = TOP + int(ROW_H * idx)

        label_box = slide.shapes.add_textbox(LEFT, y, LABEL_W, ROW_H)
        label_box.text_frame.word_wrap = True
        label_p = label_box.text_frame.paragraphs[0]
        label_p.text = task.name
        label_p.font.size = Pt(11)
        label_p.font.color.rgb = RGBColor(0x2D, 0x37, 0x48)

        start_idx = month_index.get((task.start_year, task.start_month), 0)
        end_idx = month_index.get((task.end_year, task.end_month), num_months - 1)
        bar_x = LEFT + LABEL_W + int(MONTH_W * start_idx) + Inches(0.05)
        bar_w = int(MONTH_W * (end_idx - start_idx + 1)) - Inches(0.1)
        bar_y = y + Inches(0.08)
        bar_h = ROW_H - Inches(0.16)

        bar_color = CAT_COLORS.get(task.category, RGBColor(0x99, 0x99, 0x99))
        shape = slide.shapes.add_shape(1, int(bar_x), int(bar_y), int(bar_w), int(bar_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bar_color
        shape.line.fill.background()

        shape.text_frame.paragraphs[0].text = f"{end_idx - start_idx + 1}개월"
        shape.text_frame.paragraphs[0].font.size = Pt(8)
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════
# 내보내기 엔드포인트
# ══════════════════════════════════════════════

@router.post("/export")
async def export_timeline(request: ExportRequest):
    """타임라인 내보내기 (PNG/XLSX+TODO/PPTX)"""
    try:
        detail_groups = request.detail_tasks or []
        fmt = request.format

        if fmt == "png":
            data = _generate_png(request.timeline)
            filename = f"{request.timeline.title}_일정표.png"
            mime = "image/png"
        elif fmt == "xlsx":
            data = _generate_xlsx(request.timeline, detail_groups)
            filename = f"{request.timeline.title}_일정표.xlsx"
            mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif fmt == "pptx":
            data = _generate_pptx(request.timeline)
            filename = f"{request.timeline.title}_일정표.pptx"
            mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            raise HTTPException(status_code=400, detail="지원하지 않는 형식")

        return {
            "success": True,
            "data": base64.b64encode(data).decode("utf-8"),
            "filename": filename,
            "mime_type": mime,
            "format": fmt,
        }

    except Exception as e:
        logger.error(f"내보내기 실패: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ══════════════════════════════════════════════
# 헬스체크
# ══════════════════════════════════════════════

@router.get("/status")
async def timeline_status():
    return {
        "status": "ok",
        "version": "v5",
        "features": {
            "auto_suggest": True,
            "detail_tasks_with_duration": True,
            "law_chatbot_pinpoint_queries": True,
            "execute_phase_2pass": True,
            "export_png": True,
            "export_xlsx_with_todo": True,
            "export_pptx": True,
        },
        "project_types": len(PROJECT_TYPES),
        "contract_types": len(CONTRACT_TYPES),
        "categories": len(CATEGORIES),
    }